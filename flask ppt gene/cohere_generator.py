
import os
import re
from pptx import Presentation
from crawlers.icrawlercrawler import ICrawlerCrawler
from apis.cohere_api import CohereAPIClient
from utils import get_config

def generate_ppt(topic, model_name, num_slides, theme_choice):
    config = get_config()
    api_key = config.get('api_key')

    legal_topic = re.sub(r'[^\w\s-]', '', topic).strip().replace(' ', '_')
    save_dir = os.path.join(config["save_location"], legal_topic)
    os.makedirs(save_dir, exist_ok=True)

    # Theme selection
    theme_map = {
        "light": "theme0.pptx",
        "dark": "theme1.pptx",
        "aesthetic": "theme2.pptx"
    }
    theme_file = theme_map.get(theme_choice, "theme0.pptx")

    ppt = Presentation(theme_file)

    # Updated prompt to request 3 bullet points minimum
    final_prompt = f"""
    Create an outline for a {num_slides}-slide presentation on '{topic}'.
    Slide types:
    - Title Slide: [L_TS] with [TITLE][/TITLE] and [SUBTITLE][/SUBTITLE]
    - Content Slide: [L_CS] with [TITLE][/TITLE] and [CONTENT][/CONTENT] (minimum 3 bullet points)
    - Image Slide: [L_IS] with [TITLE][/TITLE], [CONTENT][/CONTENT], and [IMAGE][/IMAGE]
    - Thanks Slide: [L_THS] with [TITLE][/TITLE]
    Separate slides using [SLIDEBREAK].

    After the content slides:
    1. Add a 'Summary' slide with a brief overview.
    2. Add a 'Next Steps' slide listing 3 clear actions.

    Use bullet points wherever possible.
    """

    api_client = CohereAPIClient(api_key, model_name)
    presentation_content = api_client.generate(final_prompt)

    def delete_all_slides():
        for i in range(len(ppt.slides) - 1, -1, -1):
            r_id = ppt.slides._sldIdLst[i].rId
            ppt.part.drop_rel(r_id)
            del ppt.slides._sldIdLst[i]

    def create_title_slide(title, subtitle):
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def create_content_slide(title, content):
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = content

    def create_image_slide(title, content, image_query):
        slide = ppt.slides.add_slide(ppt.slide_layouts[8])
        slide.shapes.title.text = title
        slide.placeholders[2].text = content
        crawler = ICrawlerCrawler(browser="google")
        image_name = crawler.get_image(image_query, save_dir)

        if not image_name:
            print(f"[⚠️] No image found for '{image_query}'. Skipping image.")
            return

        img_path = os.path.join(save_dir, image_name)
        if not os.path.exists(img_path):
            print(f"[❌] Image not found at path: {img_path}")
            return

        try:
            slide.shapes.add_picture(
                img_path,
                slide.placeholders[1].left,
                slide.placeholders[1].top,
                slide.placeholders[1].width,
                slide.placeholders[1].height
            )
        except Exception as e:
            print(f"[❌] Failed to add picture: {e}")

    def find_text_between(text, start, end):
        results = []
        start_pos = text.find(start)
        end_pos = text.find(end)
        while start_pos > -1 and end_pos > -1:
            results.append(text[start_pos + len(start):end_pos])
            start_pos = text.find(start, end_pos + len(end))
            end_pos = text.find(end, start_pos)
        return results[0] if results else ""

    def parse_response(reply):
        print("==== Cohere Reply Start ====")
        print(reply)
        print("==== Cohere Reply End ====")

        slides = reply.split("[SLIDEBREAK]")
        for slide in slides:
            if "[L_TS]" in slide:
                create_title_slide(
                    find_text_between(slide, "[TITLE]", "[/TITLE]"),
                    find_text_between(slide, "[SUBTITLE]", "[/SUBTITLE]")
                )
            elif "[L_CS]" in slide:
                create_content_slide(
                    find_text_between(slide, "[TITLE]", "[/TITLE]"),
                    find_text_between(slide, "[CONTENT]", "[/CONTENT]")
                )
            elif "[L_IS]" in slide:
                create_image_slide(
                    find_text_between(slide, "[TITLE]", "[/TITLE]"),
                    find_text_between(slide, "[CONTENT]", "[/CONTENT]"),
                    find_text_between(slide, "[IMAGE]", "[/IMAGE]")
                )
            elif "[L_THS]" in slide:
                create_content_slide(
                    find_text_between(slide, "[TITLE]", "[/TITLE]"),
                    ""
                )

    delete_all_slides()
    parse_response(presentation_content)

    # Add Summary slide
    summary_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    summary_slide.shapes.title.text = "Summary"
    summary_slide.placeholders[1].text = (
        f"This presentation provided an overview of '{topic}', "
        "highlighting the most important points and insights."
    )

     # Add Next Steps slide
    next_steps_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    next_steps_slide.shapes.title.text = "Next Steps"
    next_steps_slide.placeholders[1].text = (
        "1. Review and customize the slides.\n"
        "2. Add supporting data or visuals.\n"
        "3. Rehearse your delivery."
    )

    # Add Thank You slide
    thank_you_slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    thank_you_slide.shapes.title.text = "Thank You!" 

    # Save the final presentation
    if len(ppt.slides) > 0:
        title = ppt.slides[0].shapes.title.text.replace(":", "")
    else:
        title = "Presentation"

    output_path = os.path.join(save_dir, f"{title}.pptx")
    ppt.save(output_path)

    return output_path
